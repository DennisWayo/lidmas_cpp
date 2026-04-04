#!/usr/bin/env python3
"""Generate a simple talk deck for paper_02."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


TITLE_COLOR = RGBColor(21, 66, 129)
TEXT_COLOR = RGBColor(40, 40, 40)
ACCENT_COLOR = RGBColor(45, 125, 190)


def _set_bg(slide) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)


def _add_top_bar(slide) -> None:
    bar = slide.shapes.add_shape(
        autoshape_type_id=1,  # rectangle
        left=Inches(0.0),
        top=Inches(0.0),
        width=Inches(13.333),
        height=Inches(0.2),
    )
    fill = bar.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()


def _set_title(shape, text: str) -> None:
    shape.text = text
    p = shape.text_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR


def _set_subtitle(shape, text: str) -> None:
    shape.text = text
    p = shape.text_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_COLOR


def _set_notes(slide, notes: str) -> None:
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = notes
    for p in notes_frame.paragraphs:
        p.font.size = Pt(14)


def _title_slide(prs: Presentation, title: str, subtitle: str, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _set_bg(slide)
    _add_top_bar(slide)
    _set_title(slide.shapes.title, title)
    _set_subtitle(slide.placeholders[1], subtitle)
    _set_notes(slide, notes)


def _bullet_slide(prs: Presentation, title: str, bullets: list[str], notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_bg(slide)
    _add_top_bar(slide)
    _set_title(slide.shapes.title, title)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_COLOR
    _set_notes(slide, notes)


def _two_col_slide(prs: Presentation, title: str, left_lines: list[str], right_lines: list[str], notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_bg(slide)
    _add_top_bar(slide)

    title_shape = slide.shapes.title
    title_shape.text = title
    p = title_shape.text_frame.paragraphs[0]
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.2), Inches(5.4))
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.4))

    left_frame = left_box.text_frame
    right_frame = right_box.text_frame
    left_frame.clear()
    right_frame.clear()

    for i, line in enumerate(left_lines):
        p = left_frame.paragraphs[0] if i == 0 else left_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(22 if i == 0 else 20)
        p.font.color.rgb = TEXT_COLOR
        if i == 0:
            p.font.bold = True

    for i, line in enumerate(right_lines):
        p = right_frame.paragraphs[0] if i == 0 else right_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(22 if i == 0 else 20)
        p.font.color.rgb = TEXT_COLOR
        if i == 0:
            p.font.bold = True
    _set_notes(slide, notes)


def _closing_slide(prs: Presentation, title: str, subtitle: str, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_bg(slide)
    _add_top_bar(slide)

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(11.8), Inches(1.8))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER

    sbox = slide.shapes.add_textbox(Inches(1.2), Inches(4.3), Inches(11.0), Inches(1.2))
    sf = sbox.text_frame
    sf.clear()
    ps = sf.paragraphs[0]
    ps.text = subtitle
    ps.font.size = Pt(26)
    ps.font.color.rgb = TEXT_COLOR
    ps.alignment = PP_ALIGN.CENTER
    _set_notes(slide, notes)


def build_deck(out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _title_slide(
        prs,
        "Decoder Dependence in Surface-Code Threshold Estimation",
        "paper_02 talk deck | Dennis Delali Kwesi Wayo",
        "Open with one sentence: threshold claims depend on decoder and estimator, not only code family.",
    )

    _bullet_slide(
        prs,
        "Talk Roadmap",
        [
            "1) Why this study matters",
            "2) What we changed in methodology",
            "3) Key results and practical meaning",
            "4) Clear take-home messages",
        ],
        "Keep this short, 20-30 seconds.",
    )

    _bullet_slide(
        prs,
        "Problem and Motivation",
        [
            "Single threshold numbers can be misleading",
            "Decoder choice + estimator choice shift conclusions",
            "Need fair comparisons under matched controls",
            "Goal: reproducible, hardware-facing benchmarking",
        ],
        "Say: we focus on fair comparison first, then threshold interpretation.",
    )

    _bullet_slide(
        prs,
        "Methodology at a Glance",
        [
            "Two regimes: Pauli baseline and native GKP digitization",
            "Four decoders: MWPM, UF, BP, Neural-MWPM",
            "Runs 06-14 under matched sweep grids and seed policy",
            "LER + confidence intervals + diagnostics + bootstrap analyses",
        ],
        "Emphasize matched setup and reproducibility scripts.",
    )

    _two_col_slide(
        prs,
        "Core Experiment Setup",
        [
            "Noise and Distances",
            "Distance set: d = {3, 5, 7}",
            "Main fixed-distance tradeoff: d = 5",
            "Native GKP sigma window: 0.05 to 0.35",
            "Dense critical window: 0.08 to 0.24",
        ],
        [
            "Evaluation Protocol",
            "Matched trial budgets within each block",
            "Deterministic seeding for reruns",
            "Crossing bootstrap + rank stability",
            "Serial vs threaded fidelity checks",
        ],
        "Use this slide to show why comparisons are fair.",
    )

    _two_col_slide(
        prs,
        "Result 1: Runtime-Accuracy Tradeoff (d = 5, sigma = 0.20)",
        [
            "Pareto Frontier",
            "UF: 1.332 s, LER 0.2303",
            "MWPM: 1.341 s, LER 0.2273",
            "Neural-MWPM: 1.396 s, LER 0.3730",
            "BP: 7.640 s, LER 0.6107",
        ],
        [
            "Interpretation",
            "MWPM and UF are the practical top pair",
            "Neural-MWPM is intermediate",
            "BP is dominated in this tested regime",
            "Use MWPM/UF as baseline choices",
        ],
        "Read the numbers slowly. This is the main practical slide.",
    )

    _bullet_slide(
        prs,
        "Result 2: Crossing Stability",
        [
            "Only MWPM produced valid crossing bootstrap distributions",
            "MWPM medians: sigma*(3,5)=0.10, sigma*(5,7)=0.1375",
            "UF, BP, Neural-MWPM had no valid crossing samples in that pass",
            "Threshold scalar is estimator-conditional in this regime",
        ],
        "Say clearly: ordering is stronger evidence than one scalar crossing value.",
    )

    _bullet_slide(
        prs,
        "Result 3: Distance-Gain and Noise Sensitivity",
        [
            "Distance-gain ratios stayed below 1 across sampled sigma",
            "Larger distance did not yet lower LER in this window",
            "Measurement noise is dominant for MWPM and UF",
            "Hardware priority: improve measurement channel first",
        ],
        "Keep language simple: below-1 gain means no distance benefit in current window.",
    )

    _bullet_slide(
        prs,
        "Result 4: Ranking Robustness and Effect Sizes",
        [
            "Stable ordering across sigma: MWPM/UF at top, Neural-MWPM, then BP",
            "BP stays rank 4 and Neural-MWPM stays rank 3",
            "MWPM-UF gap is small: delta = -0.00383 (CI includes 0)",
            "MWPM strongly outperforms BP and Neural-MWPM",
        ],
        "Stress that MWPM and UF are near-tied at top in this setup.",
    )

    _bullet_slide(
        prs,
        "Result 5: Parallelization Fidelity",
        [
            "Threaded sampling improved throughput with close statistical agreement",
            "Pauli speedup: 1.34x, mean |delta LER| = 6.07e-3",
            "Native GKP speedup: 1.94x, mean |delta LER| = 5.20e-3",
            "Conclusion: faster sweeps without changing headline conclusions",
        ],
        "This slide supports practical scalability.",
    )

    _bullet_slide(
        prs,
        "Dense Critical Window Scan",
        [
            "Dense sigma scan: 0.08 to 0.24 with step 0.01",
            "All decoders returned NaN crossing entries in that estimator pass",
            "Confirms threshold localization is window-sensitive here",
            "Report threshold claims with estimator context and uncertainty",
        ],
        "Repeat: NaN crossing is an informative result, not a failure.",
    )

    _bullet_slide(
        prs,
        "Main Take-Home Messages",
        [
            "Decoder and estimator jointly determine threshold conclusions",
            "MWPM and UF provide strongest practical tradeoff in this study",
            "Ordering claims are robust; scalar threshold claims are conditional",
            "Reproducible scripts + fidelity checks are essential for fair benchmarking",
        ],
        "Pause after each bullet. This is your summary slide.",
    )

    _closing_slide(
        prs,
        "Thank You",
        "Questions and discussion",
        "If asked for one-line summary: robust decoder ordering, conditional threshold scalars.",
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    out_path = root / "paper_02_talk.pptx"
    build_deck(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    main()
